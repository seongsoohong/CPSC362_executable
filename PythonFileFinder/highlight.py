from docx import Document
from docx.enum.text import WD_COLOR_INDEX
from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
from pptx.dml.color import RGBColor
import openpyxl
from openpyxl import Workbook
from openpyxl import load_workbook
from openpyxl.styles import PatternFill, Border, Side, Alignment, Protection, Font
import os
from pynput.keyboard import Key, Controller
import time
import PyPDF2

def search_file(path, target):
    print(path)

    if '.txt' in path:
        with open(path) as f:
            if target in f.read():
                return True
            else: return False

    if '.docx' in path and '~$mpleword' not in path:
       doc = Document(path)
       print("2")
       for paragraph in doc.paragraphs:
           if target in paragraph.text:
               for run in paragraph.runs:
                   if target in run.text:
                        print("found")
                        return True
    if ".pdf" in path:
        if target == "":
            return True
        pdfFile = open(path, 'rb')
        pdfReader = PyPDF2.PdfFileReader(pdfFile)

        for index in range(pdfReader.numPages):
            pageObj = pdfReader.getPage(index)
            print(pageObj.extractText())
            if target in pageObj.extractText():
                print("fileEX")
                return True
        pdfFile.close()
        return False

    if '.pptx' in path:
        print("opening pptx")
        if target == "":
            return True
        prs = Presentation(path)
        text_runs = []
        for slides in prs.slides:
            for shape in slides.shapes:
                if not shape.has_text_frame:
                    continue
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        if target in run.text:
                            if run.text == target:
                                return True
    if '.xlsx' in path:
        print("opening xlsx")
        wb = load_workbook(path)
        for ws in wb.worksheets:
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value == None:
                        continue
                    if type(cell.value) != str :
                        temp = str(cell.value)
                        if target in temp:
                            return True
                        continue
                    if target in cell.value:
                        return True
    print("returning False")
    return False

def open_file(filename, target):
    if '.docx' in filename:
       print(filename)
       doc = Document(filename)
       for paragraph in doc.paragraphs:
           if target in paragraph.text:
               for run in paragraph.runs:

                   if target in run.text:

                       if run.text == target:
                           run.font.highlight_color = WD_COLOR_INDEX.YELLOW
                           continue
                       x = run.text.split(target)
                       print(run.text)
                       print(x)
                       print(run.text)
                       run.clear()

                       #run.add_text(x[0])
                       #temp = paragraph.add_run("", run.style)
                       #temp.add_text(target)
                       #temp.font.highlight_color = WD_COLOR_INDEX.YELLOW
                       #temp2 = paragraph.add_run(x[1], run.style)
                       paragraph.add_run(x[0], run.style)
                       temp1 = paragraph.add_run(target, run.style)
                       temp1.font.highlight_color = WD_COLOR_INDEX
                       paragraph.add_run(x[1], run.style)
                       temp2 = paragraph.add_run(target, run.style)
                       temp2.font.highlight_color = WD_COLOR_INDEX
                       paragraph.add_run(x[2], run.style)





       print("end")
       doc.save(filename)
       print("did it get saved??")
    if '.pptx' in filename:
       prs = Presentation (filename)
       text_runs= []
       for slides in prs.slides:
           for shape in slides.shapes:
               if not shape.has_text_frame:
                   continue
               text_frame = shape.text_frame
               for paragraph in text_frame.paragraphs:
                   for run in paragraph.runs:
                       if target in run.text:
                           if run.text == target:
                               run.font.fill.solid()
                               d = RGBColor(0xff, 0xff, 0x00)
                               run.font.fill.fore_color.rgb = d
                               continue
                           x = run.text.split(target);
                           run.text = x[1]
                           temp = paragraph.add_run()
                           temp.text = target
                           temp.font.bold = run.font.bold
                           temp.font.italic = run.font.italic
                           temp.font.language_id = run.font.language_id
                           temp.font.name = run.font.name
                           temp.font.size = run.font.size
                           temp.font.underline =temp.font.underline
                           temp.font.fill.solid()
                           d = RGBColor(0xff, 0xff, 0x00)
                           temp.font.fill.fore_color.rgb = d
                           temp2 = paragraph.add_run()
                           temp2.font.bold = run.font.bold
                           temp2.font.italic = run.font.italic
                           temp2.font.language_id = run.font.language_id
                           temp2.font.name = run.font.name
                           temp2.font.size = run.font.size
                           temp2.font.underline = temp.font.underline
                           temp2.text = x[1]

       prs.save(filename)
    if '.xlsx' in filename:
        wb = load_workbook(filename)
        yellowFill = PatternFill(start_color='FFFFFF00', end_color='FFFFFF00', fill_type='solid')
        for ws in wb.worksheets:
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value == None:
                        continue
                    if type(cell.value) != str :
                        temp = str(cell.value)
                        if target in temp:
                            cell.fill = yellowFill
                        continue
                    if target in cell.value:
                        cell.fill = yellowFill
        wb.save(filename)
    print("are we here?")
    os.startfile(filename)

#open_file(r'C:\Users\chris\PycharmProjects\Testing.xlsx', 'Hello')
